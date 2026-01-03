import sys
sys.path.insert(0, '/tmp/pptx_lib')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

DARK_BG = RGBColor(26, 26, 46)
CYAN = RGBColor(0, 217, 255)
PINK = RGBColor(233, 69, 96)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(150, 150, 150)

def add_dark_slide(prs):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()
    return slide

def add_text(slide, text, left, top, width, height, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.alignment = align
    p.font.color.rgb = color
    return box

def add_bullet_list(slide, items, left, top, width, height, size=16, color=WHITE):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
    return box

print("Creating slides...")

# SLIDE 1: Title
slide = add_dark_slide(prs)
add_text(slide, "BlissGene", 0, 2.5, 13.333, 1, size=72, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, "THERAPEUTICS", 0, 3.5, 13.333, 0.8, size=36, color=PINK, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, "Eliminating Human Suffering Through Precision Gene Therapy", 0, 4.5, 13.333, 0.6, size=24, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "Investor Presentation | December 2025", 0, 5.5, 13.333, 0.5, size=20, color=GRAY, align=PP_ALIGN.CENTER)
add_text(slide, "blissgene.org | info@blissgene.org", 0, 6.5, 13.333, 0.4, size=16, color=GRAY, align=PP_ALIGN.CENTER)

# SLIDE 2: The Inspiration
slide = add_dark_slide(prs)
add_text(slide, "THE INSPIRATION", 0.5, 0.5, 12, 1, size=44, color=CYAN, bold=True)
add_text(slide, '"Imagine a life where you felt minimal pain, anxiety, or fear."', 0.5, 1.5, 12, 0.8, size=28, color=CYAN)
add_text(slide, "This is the reality for Jo Cameron, a Scottish woman whose unique genetics highlights our unprecedented therapeutic strategy.", 0.5, 2.5, 12, 0.8, size=20, color=WHITE)
add_text(slide, "76", 1.5, 3.5, 2.5, 1, size=56, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, "Years pain-free", 1.5, 4.3, 2.5, 0.4, size=14, color=GRAY, align=PP_ALIGN.CENTER)
add_text(slide, "0", 5.5, 3.5, 2.5, 1, size=56, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, "Painkillers ever needed", 5.5, 4.3, 2.5, 0.4, size=14, color=GRAY, align=PP_ALIGN.CENTER)
add_text(slide, "2x", 9.5, 3.5, 2.5, 1, size=56, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, "Normal anandamide", 9.5, 4.3, 2.5, 0.4, size=14, color=GRAY, align=PP_ALIGN.CENTER)
add_text(slide, "Burns, broken bones, surgeries - all without anesthesia or pain medication.", 0.5, 5.5, 12, 0.5, size=18, color=GRAY)

# SLIDE 3: The Problem
slide = add_dark_slide(prs)
add_text(slide, "THE PROBLEM", 0.5, 0.5, 12, 1, size=44, color=CYAN, bold=True)
add_text(slide, "Massive Unmet Need", 0.5, 1.5, 5, 0.5, size=24, color=PINK, bold=True)
add_bullet_list(slide, [
    "50+ million Americans suffer from chronic pain",
    "40 million Americans have anxiety disorders",
    "17.3 million have Major Depressive Disorder",
    "41% of population affected by at least one"
], 0.5, 2.2, 5.5, 3)
add_text(slide, "Current Treatments Fail", 7, 1.5, 5, 0.5, size=24, color=PINK, bold=True)
add_bullet_list(slide, [
    "Trial-and-error medication approaches",
    "Opioid crisis with addiction risk",
    "Significant side effects",
    "Limited efficacy for many patients",
    "No treatments targeting root cause"
], 7, 2.2, 5.5, 3)
add_text(slide, "$500B+", 4, 5.5, 5, 0.8, size=48, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, "US chronic pain market", 4, 6.2, 5, 0.4, size=16, color=GRAY, align=PP_ALIGN.CENTER)

# SLIDE 4: The Science
slide = add_dark_slide(prs)
add_text(slide, "THE SCIENCE", 0.5, 0.5, 12, 1, size=44, color=CYAN, bold=True)
add_text(slide, "2023 Breakthrough: UCL Study Published in Brain Journal", 0.5, 1.3, 12, 0.5, size=22, color=PINK)
add_text(slide, "The Discovery", 0.5, 2, 5, 0.4, size=20, color=CYAN, bold=True)
add_bullet_list(slide, [
    "FAAH enzyme breaks down endocannabinoids",
    "FAAH-OUT is a lncRNA regulating FAAH",
    "Jo Cameron has 8.1 kb microdeletion",
    "Reduced FAAH = elevated anandamide"
], 0.5, 2.5, 5.5, 2.5, size=14)
add_text(slide, "Why FAAH-OUT?", 7, 2, 5, 0.4, size=20, color=CYAN, bold=True)
add_bullet_list(slide, [
    "FAAH inhibitors failed clinical trials",
    "FAAH-OUT is human-validated target",
    "Affects 1,145 genes comprehensively",
    "Includes WNT16 (healing) & BDNF (mood)"
], 7, 2.5, 5.5, 2.5, size=14)
add_text(slide, '"Our findings validate FAAH-OUT regulation of FAAH as a new route to develop pain treatments."', 0.5, 5.3, 12, 0.8, size=16, color=GRAY)
add_text(slide, "- Mikaeili et al., Brain, 2023 (DOI: 10.1093/brain/awad098)", 0.5, 6, 12, 0.4, size=12, color=GRAY)

# SLIDE 5: Our Solution
slide = add_dark_slide(prs)
add_text(slide, "OUR SOLUTION", 0.5, 0.5, 12, 1, size=44, color=CYAN, bold=True)
add_text(slide, "Targeted FAAH-OUT Knockdown Therapy", 0.5, 1.3, 12, 0.5, size=22, color=PINK)
add_text(slide, "siRNA Therapy", 0.5, 2.2, 3.5, 0.4, size=20, color=CYAN, bold=True)
add_text(slide, "Temporary knockdown\nWeeks to months (reversible)", 0.5, 2.7, 3.5, 0.8, size=14, color=WHITE)
add_text(slide, "CRISPR-Cas", 5, 2.2, 3.5, 0.4, size=20, color=CYAN, bold=True)
add_text(slide, "Permanent deletion\nMimics Jo Cameron's mutation", 5, 2.7, 3.5, 0.8, size=14, color=WHITE)
add_text(slide, "LNP Delivery", 9.5, 2.2, 3.5, 0.4, size=20, color=CYAN, bold=True)
add_text(slide, "Crosses blood-brain barrier\nProven mRNA vaccine tech", 9.5, 2.7, 3.5, 0.8, size=14, color=WHITE)
add_text(slide, "UNIQUE ADVANTAGE", 0.5, 4.5, 12, 0.5, size=20, color=PINK, bold=True)
add_text(slide, "Patients can try temporary therapy first before committing to permanent CRISPR - unprecedented patient choice!", 0.5, 5, 12, 0.6, size=18, color=WHITE)

# SLIDE 6: Market Opportunity
slide = add_dark_slide(prs)
add_text(slide, "MARKET OPPORTUNITY", 0.5, 0.5, 12, 1, size=44, color=CYAN, bold=True)
add_text(slide, "$350B", 1.5, 2, 3, 0.8, size=48, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, "TAM", 1.5, 2.7, 3, 0.4, size=20, color=PINK, align=PP_ALIGN.CENTER)
add_text(slide, "Pain + Anxiety + MDD\n67M patients @ $3-5k", 1.5, 3.2, 3, 0.6, size=12, color=GRAY, align=PP_ALIGN.CENTER)
add_text(slide, "$250B", 5.2, 2, 3, 0.8, size=48, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, "SAM", 5.2, 2.7, 3, 0.4, size=20, color=PINK, align=PP_ALIGN.CENTER)
add_text(slide, "Chronic Pain Only\n50M patients @ $5k", 5.2, 3.2, 3, 0.6, size=12, color=GRAY, align=PP_ALIGN.CENTER)
add_text(slide, "$50B", 8.9, 2, 3, 0.8, size=48, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, "SOM", 8.9, 2.7, 3, 0.4, size=20, color=PINK, align=PP_ALIGN.CENTER)
add_text(slide, "Neuropathic Pain\n10M patients @ $5k", 8.9, 3.2, 3, 0.6, size=12, color=GRAY, align=PP_ALIGN.CENTER)
add_text(slide, "Competitive Landscape", 0.5, 4.5, 12, 0.5, size=20, color=PINK, bold=True)
add_text(slide, "No other company is targeting FAAH-OUT specifically with both temporary and permanent solutions.", 0.5, 5, 12, 0.5, size=16, color=WHITE)
add_bullet_list(slide, [
    "Opioids/NSAIDs - addictive, side effects",
    "FAAH inhibitors - failed trials",
    "RNA therapeutics - not targeting this pathway"
], 0.5, 5.5, 12, 1.5, size=14)

# SLIDE 7: Leadership Team
slide = add_dark_slide(prs)
add_text(slide, "LEADERSHIP TEAM", 0.5, 0.5, 12, 1, size=44, color=CYAN, bold=True)
add_text(slide, "Brandon Emerick", 1, 2, 3.5, 0.5, size=22, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, "CEO & Founder", 1, 2.5, 3.5, 0.4, size=16, color=PINK, align=PP_ALIGN.CENTER)
add_text(slide, "Cognitive Science &\nBusiness Administration", 1, 2.9, 3.5, 0.6, size=12, color=GRAY, align=PP_ALIGN.CENTER)
add_text(slide, "Qurrat Ain, PhD", 5, 2, 3.5, 0.5, size=22, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, "COO", 5, 2.5, 3.5, 0.4, size=16, color=PINK, align=PP_ALIGN.CENTER)
add_text(slide, "Research Leadership\n& Operations", 5, 2.9, 3.5, 0.6, size=12, color=GRAY, align=PP_ALIGN.CENTER)
add_text(slide, "Valerio Embrione, PhD", 9, 2, 3.5, 0.5, size=22, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, "VP & Chief Scientific Officer", 9, 2.5, 3.5, 0.4, size=16, color=PINK, align=PP_ALIGN.CENTER)
add_text(slide, "Scientific Strategy\n& Research Direction", 9, 2.9, 3.5, 0.6, size=12, color=GRAY, align=PP_ALIGN.CENTER)
add_text(slide, "Scientific Advisory Board", 0.5, 4.2, 12, 0.5, size=20, color=PINK, bold=True)
add_text(slide, "Neena Pyzocha, PhD  |  Irfan Ullah, PhD  |  Zaki Ullah, PhD  |  Mayuri Bhattacharya (Marketing)", 0.5, 4.8, 12, 0.5, size=16, color=WHITE, align=PP_ALIGN.CENTER)

# SLIDE 8: Roadmap
slide = add_dark_slide(prs)
add_text(slide, "DEVELOPMENT ROADMAP", 0.5, 0.5, 12, 1, size=44, color=CYAN, bold=True)
timeline_items = [
    ("July 2024", "Incorporated"),
    ("Q1 2025", "Seed Funding"),
    ("Q3 2025", "Preclinical"),
    ("Q3 2027", "IND-Enabling"),
    ("Q3 2029", "IND Filing"),
    ("Q3 2030", "Phase I")
]
for i, (date, event) in enumerate(timeline_items):
    x = 1 + i * 2
    color = CYAN if i == 0 else PINK
    add_text(slide, "O", x, 2.2, 0.5, 0.4, size=24, color=color, align=PP_ALIGN.CENTER)
    add_text(slide, date, x - 0.3, 2.7, 1.1, 0.4, size=12, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, event, x - 0.3, 3.1, 1.1, 0.4, size=10, color=GRAY, align=PP_ALIGN.CENTER)
add_text(slide, "Discovery & Preclinical (2-3 years)", 0.5, 4.2, 4, 0.4, size=16, color=CYAN, bold=True)
add_text(slide, "Target validation, siRNA/CRISPR optimization, in vitro/vivo studies", 0.5, 4.6, 4, 0.5, size=12, color=WHITE)
add_text(slide, "IND-Enabling (1-2 years)", 5, 4.2, 4, 0.4, size=16, color=CYAN, bold=True)
add_text(slide, "GLP toxicology, ADME studies, CMC development", 5, 4.6, 4, 0.5, size=12, color=WHITE)
add_text(slide, "Clinical Preparation (1-2 years)", 9, 4.2, 4, 0.4, size=16, color=CYAN, bold=True)
add_text(slide, "FDA submission, site selection, IRB approval", 9, 4.6, 4, 0.5, size=12, color=WHITE)

# SLIDE 9: The Ask
slide = add_dark_slide(prs)
add_text(slide, "THE ASK", 0.5, 0.5, 12, 1, size=44, color=CYAN, bold=True)
add_text(slide, "$1,000,000", 0, 2, 13.333, 1, size=64, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, "Seed Round", 0, 2.9, 13.333, 0.5, size=24, color=PINK, align=PP_ALIGN.CENTER)
add_text(slide, "Use of Funds", 0.5, 4, 6, 0.4, size=20, color=CYAN, bold=True)
add_bullet_list(slide, [
    "40% - Preclinical R&D",
    "25% - Team Expansion",
    "15% - Equipment & Lab",
    "10% - IP & Legal",
    "10% - Operations"
], 0.5, 4.5, 5, 2.5, size=14)
add_text(slide, "Exit Strategy", 7, 4, 6, 0.4, size=20, color=CYAN, bold=True)
add_text(slide, "Partnership or acquisition by major pharmaceutical company during clinical trials, following the RNA therapeutics playbook.", 7, 4.5, 5.5, 1.5, size=14, color=WHITE)

# SLIDE 10: Why Now
slide = add_dark_slide(prs)
add_text(slide, "WHY NOW?", 0.5, 0.5, 12, 1, size=44, color=CYAN, bold=True)
add_text(slide, "Scientific Timing", 0.5, 1.5, 5, 0.4, size=22, color=PINK, bold=True)
add_bullet_list(slide, [
    "2023 breakthrough - UCL published mechanism",
    "FAAH-OUT validated as therapeutic target",
    "RNA/CRISPR delivery technology mature",
    "Non-human primate models available"
], 0.5, 2, 5.5, 2.5)
add_text(slide, "Market Timing", 7, 1.5, 5, 0.4, size=22, color=PINK, bold=True)
add_bullet_list(slide, [
    "Opioid crisis creates urgent need",
    "RNA therapeutics gaining FDA acceptance",
    "No competitors targeting FAAH-OUT",
    "First-mover advantage available"
], 7, 2, 5.5, 2.5)
add_text(slide, '"With an investment of just $1 million and a few minor changes to our genome, we can become a pain-free and happy species."', 0.5, 5, 12, 0.8, size=20, color=CYAN)
add_text(slide, "- Brandon Emerick, CEO", 10, 5.6, 3, 0.4, size=14, color=GRAY)

# SLIDE 11: Thank You
slide = add_dark_slide(prs)
add_text(slide, "Thank You", 0, 2, 13.333, 1, size=64, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, "Join Us in Eliminating Human Suffering", 0, 3, 13.333, 0.6, size=28, color=PINK, align=PP_ALIGN.CENTER)
add_text(slide, "Website: blissgene.org", 0, 4.5, 13.333, 0.4, size=20, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "Email: info@blissgene.org", 0, 5, 13.333, 0.4, size=20, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "Phone: 860-483-1425", 0, 5.5, 13.333, 0.4, size=20, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "BlissGene Therapeutics | Incorporated July 19, 2024", 0, 6.5, 13.333, 0.4, size=16, color=GRAY, align=PP_ALIGN.CENTER)

prs.save('BlissGene_Investor_Presentation.pptx')
print("SUCCESS! PowerPoint created: BlissGene_Investor_Presentation.pptx")
print(f"Total slides: {len(prs.slides)}")
